import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import fitz  # PyMuPDF
import io
import requests

# Wide Page Configuration for Gamma-like Studio UI
st.set_page_config(page_title="Supriya ppt maker .ai", page_icon="🚀", layout="wide")

st.markdown("""
    <style>
    body { background-color: #0b0f19; color: #f1f5f9; }
    .main-header { font-size: 42px; font-weight: 800; background: linear-gradient(135deg, #6366f1, #a855f7, #ec4899); -webkit-background-clip: text; -webkit-text-fill-color: transparent; text-align: center; padding: 10px; font-family: 'Helvetica Neue', Arial, sans-serif; }
    .sub-header { font-size: 15px; color: #94a3b8; text-align: center; margin-bottom: 30px; }
    .gamma-card { background: #111827; border-radius: 16px; padding: 24px; border: 1px solid #1f2937; margin-bottom: 20px; box-shadow: 0 20px 25px -5px rgba(0,0,0,0.5); }
    </style>
    """, unsafe_allow_html=True)

st.markdown('<div class="main-header">🚀 Supriya ppt maker .ai</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-header">Premium Presentation Engine | Clean Layouts, Proper Fonts & High-Res Images</div>', unsafe_allow_html=True)

# SIDEBAR DESIGN STUDIO
st.sidebar.markdown("### 🎨 Executive Styles")
template_choice = st.sidebar.selectbox(
    "Choose Theme Matrix:",
    ["Premium Midnight Navy", "Physics Wallah Elite", "Minimalist Clean Light"]
)

st.sidebar.markdown("### 🔤 Typography Configuration")
font_family = st.sidebar.selectbox("Font Set:", ["Trebuchet MS", "Arial", "Georgia", "Verdana"])
title_size = st.sidebar.slider("Header Size (Pt):", 28, 44, 34)
body_size = st.sidebar.slider("Body Text Size (Pt):", 14, 24, 18)

# Color Scheme Mapping
if template_choice == "Premium Midnight Navy":
    bg_hex, title_hex, body_hex = "#0f172a", "#38bdf8", "#cbd5e1"
elif template_choice == "Physics Wallah Elite":
    bg_hex, title_hex, body_hex = "#000000", "#facc15", "#f8fafc"
else:
    bg_hex, title_hex, body_hex = "#ffffff", "#1e293b", "#475569"

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))

# 100% Reliable Image Engine
def fetch_verified_image(keyword):
    try:
        clean_kw = keyword.strip().lower()
        search_url = f"https://images.unsplash.com/photo-1457369804613-52c61a468e7d?auto=format&fit=crop&w=600&q=80"
        if "law" in clean_kw or "tort" in clean_kw or "legal" in clean_kw:
            search_url = "https://images.unsplash.com/photo-1589829545856-d10d557cf95f?auto=format&fit=crop&w=600&q=80"
        elif "science" in clean_kw or "biology" in clean_kw:
            search_url = "https://images.unsplash.com/photo-1532187863486-abf9d39d6618?auto=format&fit=crop&w=600&q=80"
            
        res = requests.get(search_url, timeout=8)
        if res.status_code == 200:
            return io.BytesIO(res.content)
    except:
        pass
    return None

col1, col2 = st.columns([1, 2])

with col1:
    st.markdown('<div class="gamma-card">', unsafe_allow_html=True)
    st.subheader("📁 Upload Module")
    uploaded_file = st.file_uploader("Drop your PDF notes here:", type=["pdf"])
    st.markdown('</div>', unsafe_allow_html=True)

with col2:
    st.subheader("🖥️ Live Content & Layout Customizer")
    
    if uploaded_file is None:
        st.info("Aapki PDF ka text clean karke proper heading aur points mein convert kiya jayega. Kripya file upload karein.")
    else:
        pdf_bytes = uploaded_file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        r_bg, g_bg, b_bg = hex_to_rgb(bg_hex)
        r_t, g_t, b_t = hex_to_rgb(title_hex)
        r_b, g_b, b_b = hex_to_rgb(body_hex)
        
        for index, page in enumerate(doc):
            raw_text = page.get_text()
            lines = [line.strip() for line in raw_text.split('\n') if line.strip() and "CLASSES" not in line.upper() and "MEERUT" not in line.upper() and "88514" not in line]
            
            if lines:
                # Intelligently clean titles instead of using class headers
                potential_title = lines[0]
                if len(potential_title) > 50 or "semester" in potential_title.lower() or "notes" in potential_title.lower():
                    slide_title = "Core Concept Study"
                    content_start = 0
                else:
                    slide_title = potential_title
                    content_start = 1
                
                points = []
                for l in lines[content_start:7]:
                    if l.startswith(('1.', '2.', '3.', '4.')) or "?" in l:
                        points.append(f"• {l}")
                    else:
                        if points:
                            points[-1] += f" {l}"
                        else:
                            points.append(f"• {l}")
                
                slide_content = "\n\n".join(points) if points else "Core subject details compiled for review."
                
                with st.expander(f"⚙️ Slide Layout {index + 1} Editor", expanded=(index==0)):
                    edited_title = st.text_input(f"Slide {index+1} Header:", value=slide_title, key=f"t_{index}")
                    edited_body = st.text_area(f"Slide {index+1} Points:", value=slide_content, key=f"b_{index}")
                
                # PowerPoint Engineering
                blank_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_layout)
                
                # Set Solid Modern Background
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(r_bg, g_bg, b_bg)
                
                # Formatted Title Box
                title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.2))
                tf_t = title_box.text_frame
                tf_t.word_wrap = True
                p_title = tf_t.paragraphs[0]
                p_title.text = edited_title.upper()
                p_title.font.name = font_family
                p_title.font.size = Pt(title_size)
                p_title.font.bold = True
                p_title.font.color.rgb = RGBColor(r_t, g_t, b_t)
                
                # Formatted Clean Body Box (Left Side)
                body_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.2), Inches(4.5))
                tf_body = body_box.text_frame
                tf_body.word_wrap = True
                
                for p_idx, pt in enumerate(edited_body.split('\n\n')):
                    if p_idx == 0:
                        p_body = tf_body.paragraphs[0]
                    else:
                        p_body = tf_body.add_paragraph()
                    p_body.text = pt
                    p_body.font.name = font_family
                    p_body.font.size = Pt(body_size)
                    p_body.font.color.rgb = RGBColor(r_b, g_b, b_b)
                    p_body.space_after = Pt(14)
                    p_body.line_spacing = 1.2
                
                # Image Placement (Right Side Box - Seamless without borders)
                img_stream = fetch_verified_image(edited_title)
                if img_stream:
                    slide.shapes.add_picture(img_stream, Inches(7.6), Inches(2.2), width=Inches(4.9), height=Inches(4.2))

        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        st.markdown("---")
        st.balloons()
        st.success("🎉 Gamma-Style Layout Compiled Perfectly!")
        st.download_button(
            label="📥 Download Fixed Professional Presentation",
            data=ppt_buffer,
            file_name="Supriya_Premium_AI_Presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )